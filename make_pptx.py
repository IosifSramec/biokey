from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG    = RGBColor(0x0d, 0x1a, 0x0d)
SURF  = RGBColor(0x14, 0x2a, 0x14)
SURF2 = RGBColor(0x1a, 0x35, 0x1a)
ACC   = RGBColor(0x4a, 0x9a, 0x4a)
ACC2  = RGBColor(0x2d, 0x7a, 0x2d)
LIGHT = RGBColor(0xe8, 0xf5, 0xe0)
MUTED = RGBColor(0x6a, 0xaa, 0x6a)
RED   = RGBColor(0xb9, 0x1c, 0x1c)
YELLOW= RGBColor(0xd4, 0xc5, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def t(s, text, x, y, w, h, size=20, bold=False, color=LIGHT,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = 'Courier New'


def bx(s, x, y, w, h, fill=SURF, border=ACC, bw=Pt(1.2)):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = bw
    return sh


def hline(s, y, color=ACC, thickness=0.04):
    sh = s.shapes.add_shape(1, Inches(0), Inches(y), W, Inches(thickness))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
hline(s, 0.5)
hline(s, 6.85)

t(s, '// BIOKEY SYSTEM v2.0', 0.7, 0.7, 12, 0.5, size=13, color=MUTED)
t(s, 'Your heartbeat is the key —\nat any stress level.', 0.7, 1.3, 11.5, 2.6,
  size=52, bold=True, color=LIGHT)
t(s, 'Stress-Invariant Biometric Authentication  ·  AES-256-GCM',
  0.7, 4.1, 11.5, 0.6, size=20, color=ACC)
hline(s, 4.95, SURF2, 0.06)
t(s, 'HackTM 2026   //   Defence Track',
  0.7, 5.15, 12, 0.45, size=14, color=MUTED)
t(s, 'Jigmond Alexandru   ·   Milotin Rareș   ·   Sramec Iosif   ·   Stan Cristian     //     UPT AC',
  0.7, 5.65, 12, 0.45, size=12, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
hline(s, 0.5)
t(s, '// 01  PROBLEM', 0.7, 0.62, 12, 0.4, size=12, color=MUTED)
t(s, 'Authentication breaks\nwhen it matters most.', 0.7, 1.05, 11, 1.7,
  size=40, bold=True, color=LIGHT)

cards = [
    ('STRESS',          'A soldier\'s ECG at 70 bpm\nis not the same at 160 bpm.\nThe system denies its own operator.'),
    ('CAPTURE',         'If the device is seized,\nall onboard data is immediately\naccessible. No cryptographic barrier.'),
    ('COERCION',        'Passwords and cards\ncan be taken by force.\nThe operator cannot protect them.'),
    ('SINGLE TEMPLATE', 'Every existing biometric system\nenrolls once, at rest.\nField conditions break the assumption.'),
]
for i, (title, desc) in enumerate(cards):
    cx = 0.4 + i * 3.15
    bx(s, cx, 3.0, 3.0, 2.9, fill=SURF, border=RED, bw=Pt(1.5))
    t(s, title, cx+0.15, 3.12, 2.75, 0.4, size=11, bold=True, color=RED)
    t(s, desc,  cx+0.15, 3.62, 2.75, 2.0, size=13, color=MUTED)

hline(s, 6.85)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
hline(s, 0.5)
t(s, '// 02  SOLUTION', 0.7, 0.62, 12, 0.4, size=12, color=MUTED)
t(s, 'Operator-bound encryption.\nStress-invariant.', 0.7, 1.05, 12, 1.5,
  size=38, bold=True, color=LIGHT)

points = [
    ('MULTI-TEMPLATE ENROLLMENT',
     'The operator is enrolled across 5 physiological states — rest to heavy exertion. The system recognizes them at any heart rate, automatically.'),
    ('BIOMETRIC KEY DERIVATION',
     'The AES-256 encryption key is derived directly from the cardiac signature + PIN. Without the right person, the data is cryptographically inaccessible — not just access-restricted.'),
    ('SECOND FACTOR — PIN',
     'Even with correct biometrics, a wrong PIN fails at the cryptographic level. Captured biometric data alone is not enough.'),
]
for i, (title, desc) in enumerate(points):
    cy = 2.75 + i * 1.22
    bx(s, 0.4, cy, 12.5, 1.08, fill=SURF, border=ACC, bw=Pt(1.5))
    t(s, title, 0.6, cy+0.1,  3.5,  0.35, size=11, bold=True, color=ACC)
    t(s, desc,  4.3, cy+0.1, 8.4, 0.85, size=14, color=LIGHT)

hline(s, 6.85)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — BUSINESS
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
hline(s, 0.5)
t(s, '// 03  IMPACT & SCALABILITY', 0.7, 0.62, 12, 0.4, size=12, color=MUTED)
t(s, 'Any role where the right person\nmust be at the controls — right now.',
  0.7, 1.05, 11, 1.4, size=34, bold=True, color=LIGHT)

# Left column — markets
bx(s, 0.4, 2.7, 5.9, 3.8, fill=SURF)
t(s, 'DEFENCE', 0.6, 2.82, 5.6, 0.4, size=12, bold=True, color=RED)
for i, item in enumerate([
    'Drone operator authentication',
    'Forward operating base access',
    'Classified device protection',
    'Two-man rule implementation',
]):
    t(s, f'→  {item}', 0.6, 3.38+i*0.55, 5.6, 0.5, size=15, color=LIGHT)

bx(s, 6.9, 2.7, 5.9, 3.8, fill=SURF, border=ACC, bw=Pt(1.5))
t(s, 'CIVIL & INDUSTRIAL', 7.1, 2.82, 5.6, 0.4, size=12, bold=True, color=ACC)
for i, item in enumerate([
    'Critical infrastructure access',
    'Air traffic controller verification',
    'Nuclear plant operator auth',
    'Surgeon / operating room control',
]):
    t(s, f'→  {item}', 7.1, 3.38+i*0.55, 5.6, 0.5, size=15, color=LIGHT)

t(s, 'Cryptographic layer is production-ready.  Path to deployment: TPM/HSM integration — no fundamental blockers.',
  0.7, 6.58, 12, 0.35, size=12, italic=True, color=MUTED)
hline(s, 6.85)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save('BioKey_Presentation.pptx')
print("Saved: BioKey_Presentation.pptx  (4 slides)")
